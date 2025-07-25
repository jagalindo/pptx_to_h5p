"""Simple PowerPoint (PPTX) to H5P Course Presentation converter.

This script extracts images, text (with basic formatting), simple shapes and
media from a PowerPoint file and builds the folder structure expected for an
H5P Course Presentation.  Passing ``--pack`` will also copy the standard H5P
libraries from the ``jagalindo/h5p-cli`` Docker image and zip the directory
into a ``.h5p`` file.
"""

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
import os
import json
import subprocess
import zipfile


def _docker_cat(path):
    """Return the contents of ``path`` inside the ``jagalindo/h5p-cli`` image."""
    result = subprocess.run(
        [
            "docker",
            "run",
            "--rm",
            "jagalindo/h5p-cli",
            "cat",
            path,
        ],
        capture_output=True,
        text=True,
        check=True,
    )
    return result.stdout


def _collect_dependencies(machine_name, major, minor, seen):
    """Recursively collect dependency directories for a library."""
    key = f"{machine_name}-{major}.{minor}"
    if key in seen:
        return
    seen.add(key)

    # Read the library.json from the Docker image to find further dependencies
    lib_path = f"/usr/local/lib/h5p/{key}/library.json"
    try:
        data = json.loads(_docker_cat(lib_path))
    except subprocess.CalledProcessError:
        # If the library isn't found just record it and continue
        return

    for dep in data.get("preloadedDependencies", []):
        _collect_dependencies(
            dep.get("machineName"),
            dep.get("majorVersion"),
            dep.get("minorVersion"),
            seen,
        )


def copy_extensions(target_dir):
    """Copy only the required H5P libraries into ``target_dir``."""

    h5p_json_path = os.path.join(target_dir, "h5p.json")
    with open(h5p_json_path, "r", encoding="utf-8") as f:
        h5p_def = json.load(f)

    libs = set()
    for dep in h5p_def.get("preloadedDependencies", []):
        _collect_dependencies(
            dep.get("machineName"),
            dep.get("majorVersion"),
            dep.get("minorVersion"),
            libs,
        )

    if not libs:
        return

    # Copy the collected libraries from the Docker image
    abs_target = os.path.abspath(target_dir)
    for lib in libs:
        subprocess.run(
            [
                "docker",
                "run",
                "--rm",
                "-v",
                f"{abs_target}:/data",
                "jagalindo/h5p-cli",
                "sh",
                "-c",
                f"mkdir -p /data/.h5p/libraries && cp -r /usr/local/lib/h5p/{lib} /data/.h5p/libraries/",
            ],
            check=True,
        )

def emu_to_px(emu):
    """Convert EMU (English Metric Unit) to pixels (assuming 96 DPI)."""
    if emu is None:
        return 0
    return int(emu / 9525)  # 1 pixel = 9525 EMUs at 96 DPI


def create_h5p_archive(source_dir, archive_path=None):
    """Create a .h5p file from ``source_dir`` following ``h5p-cli pack`` semantics."""
    if archive_path is None:
        archive_path = os.path.abspath(source_dir) + ".h5p"
    elif not archive_path.endswith(".h5p"):
        archive_path += ".h5p"

    library_root = os.path.join(source_dir, ".h5p", "libraries")

    with zipfile.ZipFile(archive_path, "w", zipfile.ZIP_DEFLATED) as zf:
        for root, _, files in os.walk(source_dir):
            # Skip the internal .h5p directory altogether
            if os.path.relpath(root, source_dir).startswith(".h5p"):
                continue
            for fname in files:
                fpath = os.path.join(root, fname)
                arcname = os.path.relpath(fpath, source_dir)
                zf.write(fpath, arcname)

        # Libraries are stored under .h5p/libraries but must appear at the
        # archive root.  Only copy the libraries themselves.
        if os.path.isdir(library_root):
            for lib in os.listdir(library_root):
                lib_path = os.path.join(library_root, lib)
                for root, _, files in os.walk(lib_path):
                    for fname in files:
                        fpath = os.path.join(root, fname)
                        arcname = os.path.relpath(fpath, library_root)
                        zf.write(fpath, arcname)
    return archive_path

def convert_pptx_to_h5p(input_pptx, output_dir='h5p_content', pack=False):
    """
    Converts a PPTX file into an H5P Course Presentation package structure.

    Parameters:
      input_pptx -- path to the ``.pptx`` file.
      output_dir -- destination folder for the generated H5P directory tree.
      pack -- when ``True`` copy the H5P libraries using Docker and create
              a ``.h5p`` archive by zipping the directory.

    The resulting folder contains ``h5p.json`` plus a ``content`` directory
    with ``content.json`` and copied media assets.  Images are stored in
    ``images/`` and audio or video in ``media/``.
    """
    try:
        prs = Presentation(input_pptx)
    except Exception as exc:
        raise RuntimeError(f"Unable to open PPTX file: {exc}")
    content_dir = os.path.join(output_dir, 'content')
    images_dir = os.path.join(content_dir, 'images')
    media_dir = os.path.join(content_dir, 'media')
    os.makedirs(images_dir, exist_ok=True)
    os.makedirs(media_dir, exist_ok=True)

    slides = []
    for idx, slide in enumerate(prs.slides, start=1):
        slide_dict = {'elements': []}
        for shape_idx, shape in enumerate(slide.shapes, start=1):
            left  = emu_to_px(shape.left)
            top   = emu_to_px(shape.top)
            width = emu_to_px(shape.width)
            height= emu_to_px(shape.height)

            # Handle images
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                image = shape.image
                ext   = image.ext
                blob  = image.blob
                img_filename = f'images/slide{idx}_img{shape_idx}.{ext}'
                img_path     = os.path.join(content_dir, img_filename)
                with open(img_path, 'wb') as img_file:
                    img_file.write(blob)
                slide_dict['elements'].append({
                    'type':   'image',
                    'path':   img_filename,
                    'x':      left,
                    'y':      top,
                    'width':  width,
                    'height': height
                })

            # Handle text
            elif shape.has_text_frame:
                paragraphs = []
                for para in shape.text_frame.paragraphs:
                    runs = []
                    for run in para.runs:
                        style = {}
                        if run.font.size:
                            style['size'] = run.font.size.pt
                        color = None
                        if run.font.color:
                            try:
                                color = run.font.color.rgb
                            except AttributeError:
                                color = None
                        if color:
                            style['color'] = str(color)
                        runs.append({'text': run.text, 'style': style})
                    paragraphs.append({'runs': runs})
                text = "\n".join(''.join(r['text'] for r in p['runs']) for p in paragraphs).strip()
                if text:
                    slide_dict['elements'].append({
                        'type':   'text',
                        'text':   text,
                        'x':      left,
                        'y':      top,
                        'width':  width,
                        'height': height,
                        'detail': paragraphs
                    })

            # Handle basic shapes
            elif shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE:
                fill_color = None
                if shape.fill.type:
                    try:
                        color_attr = shape.fill.fore_color.rgb
                    except AttributeError:
                        color_attr = None
                    if color_attr:
                        fill_color = str(color_attr)
                slide_dict['elements'].append({
                    'type':   'shape',
                    'x':      left,
                    'y':      top,
                    'width':  width,
                    'height': height,
                    'fill':   fill_color
                })

            # Handle media shapes
            elif shape.shape_type == MSO_SHAPE_TYPE.MEDIA:
                try:
                    media = shape.part.related_parts[shape._element.blip_rId]
                    ext = os.path.splitext(media.partname)[1].lstrip('.')
                    media_filename = f'media/slide{idx}_media{shape_idx}.{ext}'
                    media_path = os.path.join(content_dir, media_filename)
                    with open(media_path, 'wb') as mfile:
                        mfile.write(media.blob)
                    slide_dict['elements'].append({
                        'type':   'media',
                        'path':   media_filename,
                        'x':      left,
                        'y':      top,
                        'width':  width,
                        'height': height
                    })
                except Exception:
                    pass

        slides.append(slide_dict)

    # Write content.json
    content = {'slides': slides}
    try:
        with open(os.path.join(content_dir, 'content.json'), 'w', encoding='utf-8') as f:
            json.dump(content, f, ensure_ascii=False, indent=2)
    except OSError as exc:
        raise RuntimeError(f"Failed to write content.json: {exc}")

    # Write h5p.json (package definition)
    h5p_json = {
        "title": os.path.splitext(os.path.basename(input_pptx))[0],
        "mainLibrary": "H5P.CoursePresentation",
        "language": "en",
        "preloadedDependencies": [
            {"machineName": "H5P.CoursePresentation", "majorVersion": 1, "minorVersion": 23},
            {"machineName": "H5P.Text",               "majorVersion": 1, "minorVersion": 5},
            {"machineName": "H5P.Image",              "majorVersion": 1, "minorVersion": 3}
        ],
        "embedTypes": ["div"]
    }
    try:
        with open(os.path.join(output_dir, 'h5p.json'), 'w', encoding='utf-8') as f:
            json.dump(h5p_json, f, ensure_ascii=False, indent=2)
    except OSError as exc:
        raise RuntimeError(f"Failed to write h5p.json: {exc}")

    print(f"H5P package structure generated in '{output_dir}'.")
    if pack:
        try:
            copy_extensions(output_dir)
            archive = create_h5p_archive(output_dir)
            print(f"Packed into '{archive}'.")
        except Exception as exc:
            print(f"Packing failed: {exc}")
    else:
        print(
            "Run the Docker image 'jagalindo/h5p-cli' to copy H5P libraries "
            "and then zip the directory into an .h5p file. For example:" )
        abs_dir = os.path.abspath(output_dir)
        print(
            "    docker run --rm -v "
            f"{abs_dir}:/data jagalindo/h5p-cli sh -c 'mkdir -p /data/.h5p && cp -r /usr/local/lib/h5p/* /data/.h5p/'"
        )
        print(
            f"    (cd {abs_dir} && zip -r ../{os.path.basename(abs_dir)}.h5p .)"
        )

if __name__ == "__main__":
    import argparse
    parser = argparse.ArgumentParser(description="Convert PPTX to H5P Course Presentation")
    parser.add_argument("pptx_file", help="Path to the .pptx file to convert")
    parser.add_argument("-o", "--output", default="h5p_content",
                        help="Output directory for the H5P package structure")
    parser.add_argument(
        "--pack",
        action="store_true",
        help=(
            "Copy libraries from the jagalindo/h5p-cli Docker image and "
            "zip the result into a .h5p file"
        ),
    )
    args = parser.parse_args()

    convert_pptx_to_h5p(args.pptx_file, args.output, pack=args.pack)

